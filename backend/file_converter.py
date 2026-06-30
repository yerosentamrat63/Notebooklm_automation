import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from fpdf import FPDF

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt"}


def needs_conversion(filename: str) -> bool:
    ext = Path(filename).suffix.lower()
    return ext in {".pptx", ".ppt"}


def get_converted_filename(original: str) -> str:
    return Path(original).stem + ".pdf"


async def convert_pptx_to_pdf(input_path: str, output_dir: str) -> str:
    prs = Presentation(input_path)

    output_filename = get_converted_filename(os.path.basename(input_path))
    output_path = os.path.join(output_dir, output_filename)

    slide_width = prs.slide_width / 914400 * 25.4
    slide_height = prs.slide_height / 914400 * 25.4

    pdf = FPDF(unit="mm", format=(slide_width, slide_height))
    pdf.set_auto_page_break(auto=False)

    for slide in prs.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    font_size = 12
                    bold = False
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            bold = True

                    left = shape.left / 914400 * 25.4 if shape.left else 10
                    top = shape.top / 914400 * 25.4 if shape.top else 10
                    width = shape.width / 914400 * 25.4 if shape.width else slide_width - 20

                    pdf.set_xy(left, top)
                    if bold:
                        pdf.set_font("Helvetica", "B", font_size)
                    else:
                        pdf.set_font("Helvetica", "", font_size)

                    pdf.multi_cell(width, font_size * 0.35, text)

    pdf.output(output_path)
    return output_path
